"""
范文加载器
负责从 examples 文件夹中加载对应类型的范文内容
"""
import os
from typing import Optional

# 获取项目根目录
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
EXAMPLES_DIR = os.path.join(PROJECT_ROOT, "examples")

def load_example(template_type: str) -> Optional[str]:
    """
    加载指定类型的范文内容
    
    Args:
        template_type: 模板类型（如"重要会议"、"培训活动"等）
    
    Returns:
        范文文本内容，如果未找到则返回None
    """
    # 构建范文文件夹路径
    example_dir = os.path.join(EXAMPLES_DIR, template_type)
    
    if not os.path.exists(example_dir):
        return None
    
    # 查找文件夹中的文件
    files = os.listdir(example_dir)
    if not files:
        return None
    
    # 优先选择 .docx 文件，其次选择其他文件
    docx_files = [f for f in files if f.endswith('.docx')]
    if docx_files:
        example_file = os.path.join(example_dir, docx_files[0])
    else:
        example_file = os.path.join(example_dir, files[0])
    
    # 根据文件类型解析
    if example_file.endswith('.docx'):
        return _parse_docx(example_file)
    elif example_file.endswith('.pptx'):
        return _parse_pptx(example_file)
    elif example_file.endswith('.pdf'):
        return _parse_pdf(example_file)
    elif example_file.endswith('.png') or example_file.endswith('.jpg') or example_file.endswith('.jpeg'):
        return f"[图片范文: {os.path.basename(example_file)} - 请参考图片中的内容结构和风格]"
    else:
        # 尝试作为文本文件读取
        try:
            with open(example_file, 'r', encoding='utf-8') as f:
                return f.read()
        except:
            return None

def _parse_docx(file_path: str) -> str:
    """解析 Word 文档"""
    try:
        from docx import Document
        doc = Document(file_path)
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        return '\n'.join(text)
    except Exception as e:
        print(f"Error parsing docx file {file_path}: {e}")
        return None

def _parse_pptx(file_path: str) -> str:
    """解析 PowerPoint 文档"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        print(f"Error parsing pptx file {file_path}: {e}")
        return None

def _parse_pdf(file_path: str) -> str:
    """解析 PDF 文档"""
    try:
        import PyPDF2
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = []
            for page in reader.pages:
                text.append(page.extract_text())
            return '\n'.join(text)
    except Exception as e:
        print(f"Error parsing pdf file {file_path}: {e}")
        return None

def get_all_examples() -> dict:
    """
    获取所有可用的范文
    
    Returns:
        字典，key为模板类型，value为范文内容
    """
    examples = {}
    if not os.path.exists(EXAMPLES_DIR):
        return examples
    
    for template_type in os.listdir(EXAMPLES_DIR):
        template_path = os.path.join(EXAMPLES_DIR, template_type)
        if os.path.isdir(template_path):
            content = load_example(template_type)
            if content:
                examples[template_type] = content
    
    return examples
