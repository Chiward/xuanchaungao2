import os
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
import logging

logger = logging.getLogger(__name__)

def parse_docx(file_path: str) -> str:
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        return "\n".join(full_text)
    except Exception as e:
        logger.error(f"Error parsing DOCX {file_path}: {str(e)}")
        return ""

def parse_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        logger.error(f"Error parsing PPTX {file_path}: {str(e)}")
        return ""

def parse_pdf(file_path: str) -> str:
    try:
        reader = PdfReader(file_path)
        full_text = []
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                full_text.append(text)
        return "\n".join(full_text)
    except Exception as e:
        logger.error(f"Error parsing PDF {file_path}: {str(e)}")
        return ""

def parse_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".txt":
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
             with open(file_path, "r", encoding="gbk") as f:
                return f.read()
    else:
        raise ValueError(f"Unsupported file type: {ext}")
